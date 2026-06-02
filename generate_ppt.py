#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generate a PPTX for the provided POGZ presentation content.
Requires: python-pptx
Install: pip install python-pptx
Run: python generate_ppt.py
Outputs: POGZ_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

OUTPUT = "POGZ_presentation.pptx"

# Choose a Chinese-capable font available on your system, e.g. "Microsoft YaHei", "SimSun"
FONT_NAME = "Microsoft YaHei"

prs = Presentation()

def set_tf_font(paragraph, name=FONT_NAME, size=18, bold=False, color=(51,51,51)):
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)

def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title.text_frame
    title_tf.clear()
    title_tf.text = title
    title_tf.paragraphs[0].runs[0].font.size = Pt(40)
    title_tf.paragraphs[0].runs[0].font.name = FONT_NAME

    if slide.placeholders and len(slide.placeholders) > 1:
        sub_tf = slide.placeholders[1].text_frame
        sub_tf.clear()
        sub_tf.text = subtitle
        sub_tf.paragraphs[0].runs[0].font.size = Pt(18)
        sub_tf.paragraphs[0].runs[0].font.name = FONT_NAME

def add_bullet_slide(title, bullets, notes=None):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
        # style
        p.level = 0
        set_tf_font(p, size=20)
    # notes
    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes
        # style notes
    return slide

# Build slides from the provided outline
add_title_slide("POGZ 致病突变导致皮层发育受损与可逆的自闭症样表型",
                "Pathogenic POGZ mutation causes impaired cortical development and reversible autism-like phenotypes\n2025-2026 学年度第二学期《细胞遗传学》讨论课\n汇报人：姓名A / 姓名B")

# First part title
add_bullet_slide("【第一部分】 ASD相关新发突变及对患者NSCs的影响",
                 ["主讲人：姓名 A"],
                 notes="第一部分由姓名 A 主讲，聚焦 POGZ 突变在细胞与发育层面的证据。")

# Slide: 1. 研究背景
add_bullet_slide("1. 研究背景：ASD 与新发突变",
                 ["ASD 核心症状：社交沟通障碍、重复刻板行为",
                  "发病机制复杂，遗传异质性高",
                  "新发突变 (De novo mutations) 在散发性 ASD 中很重要",
                  "目标基因：POGZ，含锌指，高表达于大脑，常见 ASD DNM 基因之一"],
                 notes="建议配图：自闭症新发突变基因网络或 POGZ 结构图。")

# Slide: 2. 患者溯源与 Q1042R
add_bullet_slide("2. 患者溯源与 POGZ Q1042R 突变鉴定",
                 ["临床在散发性 ASD 队列中鉴定出一例携带该突变患者",
                  "POGZ 杂合新发错义突变 Q1042R（Q -> R）",
                  "科学问题：该突变如何影响早期大脑发育？"],
                 notes="建议配图：患者家系图及测序峰图。")

# Slide: 3. 实验模型构建
add_bullet_slide("3. 实验模型构建：体细胞 -> NSCs",
                 ["提取患者皮肤成纤维细胞",
                  "重编程为 iPSCs，再定向分化为 NSCs",
                  "对照：健康个体 NSCs，用于对比"],
                 notes="配图建议：重编程与分化流程(Fibroblast -> iPSC -> NSC)。")

# Slide: 4. 体外证明：NSCs 分化能力受损
add_bullet_slide("4. 体外证明：NSCs 分化能力显著受损",
                 ["患者来源 NSCs 诱导分化为 MAP2 阳性神经元的比例显著降低",
                  "结论：POGZ 突变干扰神经发生进程"],
                 notes="配图建议：文献中 NSCs 分化的免疫荧光对比图 (Figure 1)。")

# Slide: 5. 体内证明：皮层发育与迁移异常
add_bullet_slide("5. 体内证明：皮层发育与神经元迁移异常",
                 ["采用 E14.5 子宫内胚胎电转 (IUE) 敲低 Pogz",
                  "Pogz 敲低导致放射状迁移缺陷（神经元滞留在深层）",
                  "Rescue: 人类 WT POGZ 可挽救；Q1042R 无法挽救"],
                 notes="配图：IUE 后皮层切片荧光图，显示迁移位置差异。")

# Second part title
add_bullet_slide("【第二部分】 突变小鼠模型、行为异常与药物治疗",
                 ["主讲人：姓名 B"],
                 notes="第二部分由姓名 B 主讲，聚焦动物模型、行为与药物干预。")

# Slide: 6. 承上启下
add_bullet_slide("6. 承上启下：从细胞到动物",
                 ["细胞与胚胎证据无法评估成年行为影响",
                  "需求：建立携带真实致病突变并能成年的动物模型",
                  "方案：CRISPR/Cas9 构建点突变小鼠"],
                 notes="")

# Slide: 7. 小鼠模型建立
add_bullet_slide("7. 首个致病新发突变小鼠模型",
                 ["使用 CRISPR/Cas9 编辑",
                  "人 Q1042 对应鼠 Q1038，构建 POGZ WT/Q1038R 杂合子",
                  "该模型为首个携带患者真实致病突变的 POGZ 小鼠"],
                 notes="配图建议：CRISPR 打靶策略示意。")

# Slide: 8. 行为：母婴交流障碍 (USV)
add_bullet_slide("8. 行为学异常 (1)：母婴交流障碍",
                 ["超声波发声测试（USV）模拟婴儿交流异常",
                  "突变幼崽呼叫次数显著增加，持续时间更长"],
                 notes="图表建议：USV 呼叫次数与时长的柱状图。")

# Slide: 9. 行为：社交缺陷
add_bullet_slide("9. 行为学异常 (2)：核心社交互动缺陷",
                 ["Reciprocal social interaction 测试",
                  "突变小鼠对陌生同类的主动社交时间显著减少"],
                 notes="图表建议：社交互动时间统计与行为轨迹热图。")

# Slide: 10. 机制解析：E/I 失衡
add_bullet_slide("10. 机制解析：大脑皮层 E/I 失衡",
                 ["膜片钳记录显示皮层兴奋性神经元过度激活",
                  "mEPSC 频率与幅度增加 -> 兴奋性突触传递升高",
                  "导致网络 E/I 失衡，视为自闭症样表型基础"],
                 notes="插图建议：E/I 天平示意或膜片钳波形图。")

# Slide: 11. 药物治疗（Perampanel）可逆性
add_bullet_slide("11. 药物治疗的后续干预 (Pharmacological rescue)",
                 ["假设：抑制兴奋性传递可逆转表型",
                  "药物：吡仑帕奈 (Perampanel)，AMPA 受体拮抗剂",
                  "成年突变小鼠经亚致痫剂量给药后，皮层兴奋被抑制，社交缺陷得到显著改善"],
                 notes="给药流程与给药前后行为对比图建议。")

# Slide: 12. 总结与临床启发
add_bullet_slide("12. 总结与临床启发",
                 ["POGZ Q1042R 突变 -> 损害 NSC 分化与神经元迁移 -> 破坏 E/I 平衡 -> 自闭症样社交缺陷",
                  "关键启发：成年期通过调节突触兴奋性仍可部分逆转 ASD 核心表型，提示潜在治疗方向"],
                 notes="强调“可逆性”这一重要发现的临床意义。")

# Slide: 13. 参考文献 & 致谢
add_bullet_slide("13. 参考文献 & 致谢",
                 ["Matsumura, K., et al. (2020). Pathogenic POGZ mutation causes impaired cortical development and reversible autism-like phenotypes. Nature Communications, 11, 859.",
                  "小组成员：姓名 A（第一部分）、姓名 B（第二部分）"],
                 notes="感谢老师和同学们，欢迎提问 (Q&A)。")

# Optionally add a blank Q&A slide
add_bullet_slide("Q & A", ["感谢聆听！欢迎提问。"], notes="")

# Save presentation
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
